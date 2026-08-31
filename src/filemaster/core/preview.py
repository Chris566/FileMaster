"""W4 v2: 文件内容预览（文本/图片/PDF/Office/二进制）.

设计原则:
- **轻量**: 一切 IO 受 max_bytes 限制,绝不读全大文件
- **隔离**: 不同文件类型用不同 renderer,各 renderer 失败降级到 hex dump
- **可测**: 所有 IO 走 pathlib,renderer 接受 bytes (测试用) + Path (生产用)
- **跨平台**: 路径用 Path,文本解码 errors="replace" 兜底
- **可选依赖**: PyMuPDF / python-docx / openpyxl / python-pptx 任意缺失时,
  对应 renderer 降级为 "该类型暂不支持" + 基础信息

公开 API:
    PreviewKind: 枚举 (TEXT/IMAGE/PDF/OFFICE_DOC/OFFICE_SHEET/OFFICE_SLIDE/BINARY/UNSUPPORTED)
    PreviewContent: dataclass (kind + payload: str | bytes(QImage/None) | hex str | ...)
    FileMetadata: dataclass (size/mtime/ctime/mode/mime)
    build_preview(path, max_bytes) -> tuple[FileMetadata, PreviewContent]
    render_text(content, max_lines) -> str
    render_image(content) -> QImage | None
    render_pdf_page1(path) -> QImage | None
    render_office(path) -> str  # 第一段文本
    render_hex(content, max_bytes) -> str
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Optional


# ============================================================
# 类型
# ============================================================
class PreviewKind(str, Enum):
    """预览内容类型."""

    TEXT = "text"
    IMAGE = "image"
    PDF = "pdf"
    OFFICE_DOC = "office_doc"      # docx
    OFFICE_SHEET = "office_sheet"  # xlsx
    OFFICE_SLIDE = "office_slide"  # pptx
    BINARY = "binary"              # hex dump 兜底
    UNSUPPORTED = "unsupported"    # 该类型暂不支持 (库未装等)


@dataclass(frozen=True)
class FileMetadata:
    """文件元信息.

    跟 PreviewContent 分离: 元信息永远能取 (os.stat), 内容预览可能失败.
    """

    path: Path
    size: int
    mtime: float          # 修改时间戳 (epoch)
    ctime: float          # 创建时间戳 (epoch)
    mode: str             # 八进制权限字符串 如 "0644"
    mime: str             # best-effort MIME

    def to_dict(self) -> dict:
        return {
            "path": str(self.path),
            "size": self.size,
            "mtime": self.mtime,
            "ctime": self.ctime,
            "mode": self.mode,
            "mime": self.mime,
        }


@dataclass(frozen=True)
class PreviewContent:
    """预览内容.

    payload 含义按 kind 分:
        TEXT: str (前 N 行)
        IMAGE: QImage | None (None = 解码失败, 切 BINARY 兜底)
        PDF: QImage | None (None = PyMuPDF 缺失或渲染失败)
        OFFICE_*: str (第一段文字, "" = 库缺失)
        BINARY: str (hex dump)
        UNSUPPORTED: str (说明)
    """

    kind: PreviewKind
    payload: object
    note: str = ""        # 额外说明 (e.g. "PyMuPDF 未安装, 显示首 2KB")
    truncated: bool = False  # 是否被截断 (大文件)

    def to_dict(self) -> dict:
        # payload 不可序列化时不带它
        if isinstance(self.payload, (str, int, float, bool, list, dict)) or self.payload is None:
            return {
                "kind": self.kind.value,
                "payload": self.payload,
                "note": self.note,
                "truncated": self.truncated,
            }
        return {
            "kind": self.kind.value,
            "payload_type": type(self.payload).__name__,
            "note": self.note,
            "truncated": self.truncated,
        }


# ============================================================
# MIME 探测 (best-effort, 不依赖 mimetypes 标准库以保证 Windows 一致)
# ============================================================
_EXT_TO_MIME = {
    # 文本
    ".txt": "text/plain", ".md": "text/markdown", ".csv": "text/csv",
    ".json": "application/json", ".xml": "application/xml",
    ".yml": "text/yaml", ".yaml": "text/yaml", ".toml": "text/toml",
    ".ini": "text/plain", ".cfg": "text/plain", ".conf": "text/plain",
    ".log": "text/plain", ".py": "text/x-python", ".js": "text/javascript",
    ".ts": "text/typescript", ".html": "text/html", ".css": "text/css",
    ".sh": "text/x-shellscript", ".bat": "text/x-bat", ".ps1": "text/x-powershell",
    # 图片
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".bmp": "image/bmp", ".webp": "image/webp",
    ".tiff": "image/tiff", ".tif": "image/tiff",
    # 文档
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    # 旧 Office (不支持, 走 BINARY)
    ".doc": "application/msword", ".xls": "application/vnd.ms-excel",
    ".ppt": "application/vnd.ms-powerpoint",
    # 其它常见
    ".zip": "application/zip", ".rar": "application/x-rar", ".7z": "application/x-7z-compressed",
    ".tar": "application/x-tar", ".gz": "application/gzip",
    ".exe": "application/x-msdownload", ".dll": "application/x-msdownload",
    ".so": "application/x-sharedlib", ".dylib": "application/x-mach-binary",
}


def _guess_mime(path: Path) -> str:
    """按扩展名猜 MIME (不读文件, 走我们的轻量字典)."""
    return _EXT_TO_MIME.get(path.suffix.lower(), "application/octet-stream")


def _stat_safe(path: Path) -> os.stat_result | None:
    """跨平台 os.stat — 边界值/不存在时返 None."""
    try:
        return os.stat(path)
    except (OSError, ValueError):
        return None


def extract_metadata(path: Path) -> FileMetadata:
    """提取文件元信息. 文件不存在时返占位 (size=0/mtime=0/...)."""
    st = _stat_safe(path)
    if st is None:
        return FileMetadata(
            path=path,
            size=0,
            mtime=0.0,
            ctime=0.0,
            mode="0000",
            mime=_guess_mime(path),
        )
    return FileMetadata(
        path=path,
        size=st.st_size,
        mtime=st.st_mtime,
        ctime=st.st_ctime,
        # 跨平台: st_mode 在 Windows 是只读隐藏属性, 截取低 9 位即可
        mode=f"{st.st_mode & 0o7777:04o}",
        mime=_guess_mime(path),
    )


# ============================================================
# 文本渲染
# ============================================================
# 文本类扩展名 (按需读)
_TEXT_EXTS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".yml", ".yaml", ".toml",
    ".ini", ".cfg", ".conf", ".log", ".py", ".js", ".ts", ".html", ".css",
    ".sh", ".bat", ".ps1", ".sql", ".gitignore", ".env", ".properties",
    ".rst", ".tex",
}

# 文本魔数探测 (替代扩展名判定)
_TEXT_MAGIC_SAMPLE = 2048  # 读 2KB 来探测


def _is_likely_text(content: bytes) -> bool:
    """前 N 字节无 NULL 且大部分是 ASCII/UTF-8 可见字符 → 视为文本."""
    if not content:
        return True
    sample = content[:_TEXT_MAGIC_SAMPLE]
    # 含 NULL 字节基本是二进制
    if b"\x00" in sample:
        return False
    # 试 utf-8 解码
    try:
        sample.decode("utf-8")
        return True
    except UnicodeDecodeError:
        pass
    # 试 latin-1 (单字节全可解码, 至少不会抛)
    try:
        sample.decode("latin-1")
        return True
    except UnicodeDecodeError:
        return False
    return False


def render_text(content: bytes, max_lines: int = 200, max_bytes: int = 256 * 1024) -> PreviewContent:
    """渲染文本内容.

    Args:
        content: 完整文件 bytes (调用方控制大小)
        max_lines: 最多返回多少行
        max_bytes: 超过该大小则截断
    """
    truncated = len(content) > max_bytes
    sample = content[:max_bytes]
    # errors="replace" 兜底真编码不了的字符 (跨平台安全)
    text = sample.decode("utf-8", errors="replace")
    lines = text.splitlines()
    if len(lines) > max_lines:
        lines = lines[:max_lines]
        truncated = True
    payload = "\n".join(lines)
    note = ""
    if truncated:
        note = f"已截断: 原 {len(content)} 字节, 显示前 {len(sample)} 字节 / {len(lines)} 行"
    return PreviewContent(
        kind=PreviewKind.TEXT,
        payload=payload,
        note=note,
        truncated=truncated,
    )


# ============================================================
# 图片渲染
# ============================================================
def render_image(content: bytes) -> PreviewContent:
    """渲染图片. 失败时降级 BINARY.

    PySide6.QtGui.QImage 在主进程可用; 测试用 bytes 走 stub.
    """
    try:
        from PySide6.QtGui import QImage  # type: ignore
    except ImportError:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="PySide6 未安装, 无法渲染图片",
        )
    img = QImage()
    if not img.loadFromData(content):
        return PreviewContent(
            kind=PreviewKind.BINARY,
            payload=render_hex(content, max_bytes=512).payload,
            note="图片解码失败, 显示 hex dump",
            truncated=True,
        )
    return PreviewContent(kind=PreviewKind.IMAGE, payload=img, note="")


# ============================================================
# PDF 渲染 (取第 1 页)
# ============================================================
def render_pdf_page1(path: Path) -> PreviewContent:
    """渲染 PDF 第 1 页. PyMuPDF 缺失或失败 → UNSUPPORTED.

    返回 payload: QImage | None (None 时 note 会说明原因).
    """
    try:
        import fitz  # PyMuPDF  # type: ignore
    except ImportError:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="PyMuPDF 未安装, 无法渲染 PDF",
        )
    try:
        doc = fitz.open(str(path))
        if doc.page_count == 0:
            doc.close()
            return PreviewContent(
                kind=PreviewKind.PDF,
                payload=None,
                note="PDF 无页",
            )
        page = doc.load_page(0)
        # 1.5x 缩放 → 大约 150 DPI
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))
        # 走 QImage
        try:
            from PySide6.QtGui import QImage  # type: ignore
        except ImportError:
            doc.close()
            return PreviewContent(
                kind=PreviewKind.UNSUPPORTED,
                payload="",
                note="PySide6 未安装, PDF 已渲染为 pixmap 但无法转 QImage",
            )
        img = QImage(
            pix.samples,
            pix.width,
            pix.height,
            pix.stride,
            QImage.Format.Format_RGB888,
        )
        doc.close()
        return PreviewContent(kind=PreviewKind.PDF, payload=img, note="第 1 页 (1.5x)")
    except Exception as e:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"PDF 渲染失败: {e}",
        )


# ============================================================
# Office 渲染 (取第 1 段)
# ============================================================
def render_office(path: Path) -> PreviewContent:
    """按扩展名分发: docx / xlsx / pptx. 库缺失时降级 BINARY."""
    ext = path.suffix.lower()
    if ext == ".docx":
        return _render_docx(path)
    if ext == ".xlsx":
        return _render_xlsx(path)
    if ext == ".pptx":
        return _render_pptx(path)
    # 旧格式 (.doc/.xls/.ppt) — 不支持
    return PreviewContent(
        kind=PreviewKind.UNSUPPORTED,
        payload="",
        note=f"旧版 Office 格式 {ext} 暂不支持, 请另存为新版 (.docx/.xlsx/.pptx)",
    )


def _render_docx(path: Path) -> PreviewContent:
    try:
        from docx import Document  # python-docx  # type: ignore
    except ImportError:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="python-docx 未安装, 无法读取 .docx",
        )
    try:
        doc = Document(str(path))
        # 取前 50 段
        paras = []
        for i, p in enumerate(doc.paragraphs):
            if i >= 50:
                paras.append(f"... (还有 {len(doc.paragraphs) - 50} 段)")
                break
            txt = p.text.strip()
            if txt:
                paras.append(txt)
        if not paras:
            return PreviewContent(
                kind=PreviewKind.OFFICE_DOC,
                payload="(空文档)",
                note="无文本段落",
            )
        return PreviewContent(
            kind=PreviewKind.OFFICE_DOC,
            payload="\n".join(paras),
            note=f"Word 文档 · {len(doc.paragraphs)} 段",
        )
    except Exception as e:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"读取 .docx 失败: {e}",
        )


def _render_xlsx(path: Path) -> PreviewContent:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="openpyxl 未安装, 无法读取 .xlsx",
        )
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        out_lines = [f"📊 Workbook: {len(wb.sheetnames)} sheet(s)"]
        for sname in wb.sheetnames[:3]:  # 最多 3 个 sheet
            out_lines.append(f"\n--- {sname} ---")
            ws = wb[sname]
            for i, row in enumerate(ws.iter_rows(max_row=20, values_only=True)):
                if i >= 20:
                    out_lines.append("... (还有行)")
                    break
                # 截断每列避免超宽
                cells = [str(c)[:30] if c is not None else "" for c in row[:8]]
                out_lines.append(" | ".join(cells))
        wb.close()
        return PreviewContent(
            kind=PreviewKind.OFFICE_SHEET,
            payload="\n".join(out_lines),
            note="前 3 sheet × 前 20 行",
        )
    except Exception as e:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"读取 .xlsx 失败: {e}",
        )


def _render_pptx(path: Path) -> PreviewContent:
    try:
        from pptx import Presentation  # python-pptx  # type: ignore
    except ImportError:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="python-pptx 未安装, 无法读取 .pptx",
        )
    try:
        prs = Presentation(str(path))
        out_lines = [f"🎞️ Presentation: {len(prs.slides)} slide(s)"]
        for i, slide in enumerate(prs.slides[:5]):  # 前 5 slide
            out_lines.append(f"\n--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            out_lines.append(txt)
        return PreviewContent(
            kind=PreviewKind.OFFICE_SLIDE,
            payload="\n".join(out_lines),
            note=f"前 5 slide ({len(prs.slides)} 总)",
        )
    except Exception as e:
        return PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"读取 .pptx 失败: {e}",
        )


# ============================================================
# 二进制 hex dump
# ============================================================
def render_hex(content: bytes, max_bytes: int = 1024) -> PreviewContent:
    """Hex + ASCII 视图, 默认前 1KB."""
    sample = content[:max_bytes]
    truncated = len(content) > max_bytes

    lines = []
    for offset in range(0, len(sample), 16):
        chunk = sample[offset:offset + 16]
        hex_part = " ".join(f"{b:02x}" for b in chunk)
        ascii_part = "".join(chr(b) if 32 <= b < 127 else "." for b in chunk)
        lines.append(f"{offset:08x}  {hex_part:<48}  {ascii_part}")

    note = ""
    if truncated:
        note = f"已截断: 原 {len(content)} 字节, 显示前 {len(sample)} 字节"
    return PreviewContent(
        kind=PreviewKind.BINARY,
        payload="\n".join(lines),
        note=note,
        truncated=truncated,
    )


# ============================================================
# 顶层工厂: build_preview
# ============================================================
_TEXT_PREVIEW_MAX_BYTES = 256 * 1024  # 256KB
_HEX_PREVIEW_MAX_BYTES = 1024        # 1KB


def build_preview(path: Path, max_text_bytes: int = _TEXT_PREVIEW_MAX_BYTES) -> tuple[FileMetadata, PreviewContent]:
    """根据文件类型/内容选 renderer, 统一入口.

    Returns:
        (FileMetadata, PreviewContent) — 前者永远能取, 后者尽力.

    设计: 任何 renderer 内部失败都会被本函数 catch, 降级到 BINARY hex.
    """
    meta = extract_metadata(path)

    # 文件不存在/为空
    if meta.size == 0:
        return meta, PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note="空文件或文件不存在",
        )

    # 文件太大 (超 50MB) 直接走 hex, 不读全文
    if meta.size > 50 * 1024 * 1024:
        try:
            with open(path, "rb") as f:
                head = f.read(_HEX_PREVIEW_MAX_BYTES)
        except OSError as e:
            return meta, PreviewContent(
                kind=PreviewKind.UNSUPPORTED,
                payload="",
                note=f"读取失败: {e}",
            )
        hex_content = render_hex(head, max_bytes=_HEX_PREVIEW_MAX_BYTES)
        return meta, PreviewContent(
            kind=PreviewKind.BINARY,
            payload=hex_content.payload,
            note=f"文件 {meta.size // (1024 * 1024)} MB 过大, 仅显示头部 1KB hex",
            truncated=True,
        )

    ext = path.suffix.lower()
    mime = meta.mime

    # 按 MIME 分发
    try:
        if mime.startswith("image/"):
            with open(path, "rb") as f:
                content = f.read()
            return meta, render_image(content)
        if mime == "application/pdf":
            return meta, render_pdf_page1(path)
        if ext == ".docx":
            return meta, render_office(path)
        if ext == ".xlsx":
            return meta, render_office(path)
        if ext == ".pptx":
            return meta, render_office(path)
    except Exception as e:
        return meta, PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"渲染失败: {e}",
        )

    # 文本判定: 按扩展名 OR 按内容探测
    try:
        with open(path, "rb") as f:
            content = f.read(max_text_bytes + 1)  # 多读 1 字节判断是否截断
    except OSError as e:
        return meta, PreviewContent(
            kind=PreviewKind.UNSUPPORTED,
            payload="",
            note=f"读取失败: {e}",
        )

    if ext in _TEXT_EXTS or _is_likely_text(content):
        return meta, render_text(content, max_bytes=max_text_bytes)

    # 兜底: hex
    return meta, render_hex(content, max_bytes=_HEX_PREVIEW_MAX_BYTES)


# ============================================================
# 测试用 helper
# ============================================================
def classify_for_preview(path: Path) -> PreviewKind:
    """不读内容, 仅按扩展名判断 kind (测试 + 调度用)."""
    mime = _guess_mime(path)
    if mime.startswith("image/"):
        return PreviewKind.IMAGE
    if mime == "application/pdf":
        return PreviewKind.PDF
    ext = path.suffix.lower()
    if ext == ".docx":
        return PreviewKind.OFFICE_DOC
    if ext == ".xlsx":
        return PreviewKind.OFFICE_SHEET
    if ext == ".pptx":
        return PreviewKind.OFFICE_SLIDE
    if ext in _TEXT_EXTS:
        return PreviewKind.TEXT
    return PreviewKind.BINARY


# ============================================================
# PreviewGenerator — 给 core/__init__.py 用的类包装
# ============================================================
class PreviewGenerator:
    """文件内容预览器 (类包装, W4 v2 GUI/CLI 共用).

    Usage:
        gen = PreviewGenerator()
        meta, content = gen.generate(path)   # 等价 build_preview(path)
    """

    def __init__(self, max_text_bytes: int = _TEXT_PREVIEW_MAX_BYTES) -> None:
        self.max_text_bytes = max_text_bytes

    def generate(self, path: Path) -> tuple[FileMetadata, PreviewContent]:
        return build_preview(path, max_text_bytes=self.max_text_bytes)

    def metadata(self, path: Path) -> FileMetadata:
        return extract_metadata(path)

    def classify(self, path: Path) -> PreviewKind:
        return classify_for_preview(path)
